from typing import cast

import daiquiri
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData, _BaseChartData
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.shapetree import SlideShapes
from pptx.slide import Slide

from dbnomics_pptx_tools.pptx_copy import copy_shape_properties
from dbnomics_pptx_tools.xml_utils import remove_element

logger = daiquiri.getLogger(__name__)


def recreate_chart(chart_shape: GraphicFrame, *, chart_data: CategoryChartData, slide: Slide) -> Chart:
    chart = cast(Chart, chart_shape.chart)
    remove_element(chart_shape.element)
    new_chart_shape = cast(
        GraphicFrame,
        cast(SlideShapes, slide.shapes).add_chart(
            chart.chart_type, chart_shape.left, chart_shape.top, chart_shape.width, chart_shape.height, chart_data
        ),
    )
    copy_shape_properties(chart_shape, new_chart_shape)
    logger.debug("The chart was recreated")
    return cast(Chart, new_chart_shape.chart)


def replace_chart_data_or_recreate(chart_shape: GraphicFrame, *, chart_data: _BaseChartData, slide: Slide) -> None:
    chart = cast(Chart, chart_shape.chart)

    try:
        chart.replace_data(chart_data)
    except ValueError:
        logger.debug("Failed replacing data, recreating chart", exc_info=True)
        chart = recreate_chart(chart_shape, chart_data=chart_data, slide=slide)
    else:
        logger.debug("Chart data was replaced")
