from collections.abc import Iterable
from typing import cast

from pptx.chart.axis import AxisTitle, CategoryAxis, DateAxis, MajorGridlines, TickLabels, ValueAxis, _BaseAxis
from pptx.chart.chart import Chart, ChartTitle, _Plots
from pptx.chart.datalabel import DataLabel, DataLabels
from pptx.chart.legend import Legend
from pptx.chart.plot import _BasePlot
from pptx.chart.point import Point
from pptx.chart.series import SeriesCollection, XySeries, _BaseCategorySeries
from pptx.dml.chtfmt import ChartFormat
from pptx.dml.color import ColorFormat
from pptx.dml.fill import FillFormat
from pptx.dml.line import LineFormat
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.shapes.autoshape import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.text.text import Font, TextFrame


def copy_axis_properties(source: _BaseAxis, target: _BaseAxis) -> None:
    assert type(source) == type(target)
    copy_scalar_attributes(
        source,
        target,
        attribute_names={
            "has_major_gridlines",
            "has_minor_gridlines",
            "has_title",
            "maximum_scale",
            "minimum_scale",
            "major_tick_mark",
            "minor_tick_mark",
            "reverse_order",
            "tick_label_position",
            "visible",
        },
    )
    copy_chart_format_properties(cast(ChartFormat, source.format), cast(ChartFormat, target.format))
    copy_tick_labels_properties(cast(TickLabels, source.tick_labels), cast(TickLabels, target.tick_labels), axis=source)
    if source.has_major_gridlines:
        copy_major_gridlines_properties(
            cast(MajorGridlines, source.major_gridlines), cast(MajorGridlines, target.major_gridlines)
        )
    if source.has_minor_gridlines:
        raise NotImplementedError
    if source.has_title:
        copy_axis_title_properties(source.axis_title, target.axis_title)

    if isinstance(source, DateAxis):
        pass
    elif isinstance(source, ValueAxis):
        assert isinstance(target, ValueAxis)
        copy_value_axis_properties(source, target)
    else:
        raise NotImplementedError


def copy_axis_title_properties(source: AxisTitle, target: AxisTitle) -> None:
    copy_scalar_attributes(source, target, attribute_names={"has_text_frame"})
    copy_chart_format_properties(cast(ChartFormat, source.format), cast(ChartFormat, target.format))
    if source.has_text_frame:
        copy_text_frame_properties(source.text_frame, target.text_frame)


def copy_base_category_series_properties(source: _BaseCategorySeries, target: _BaseCategorySeries) -> None:
    # Ignore data_labels because they will be generated automatically
    copy_scalar_attributes(source, target, attribute_names={"smooth"})
    copy_chart_format_properties(cast(ChartFormat, source.format), cast(ChartFormat, target.format))


def copy_chart_format_properties(source: ChartFormat, target: ChartFormat) -> None:
    copy_fill_format_properties(cast(FillFormat, source.fill), cast(FillFormat, target.fill))
    copy_line_format_properties(cast(LineFormat, source.line), cast(LineFormat, target.line))


def copy_chart_properties(source: Chart, target: Chart) -> None:
    """Copy chart properties.

    Plots are ignored because they are set by chart data when calling add_chart.
    """
    assert source.chart_type == target.chart_type  # Set by add_chart
    copy_scalar_attributes(source, target, attribute_names={"has_legend", "has_title"})
    copy_axis_properties(source.category_axis, target.category_axis)
    copy_axis_properties(source.value_axis, target.value_axis)
    copy_font_properties(cast(Font, source.font), cast(Font, target.font))
    if source.chart_style is not None:
        target.chart_style = source.chart_style
    if source.has_legend:
        copy_legend_properties(cast(Legend, source.legend), cast(Legend, target.legend))
    if source.has_title:
        copy_chart_title_properties(source.chart_title, target.chart_title)
    copy_series_collection_properties(cast(SeriesCollection, source.series), cast(SeriesCollection, target.series))
    copy_plots_properties(cast(_Plots, source.plots), cast(_Plots, target.plots))


def copy_chart_title_properties(source: ChartTitle, target: ChartTitle) -> None:
    copy_axis_title_properties(cast(AxisTitle, source), cast(AxisTitle, target))


def copy_color_format_properties(source: ColorFormat, target: ColorFormat) -> None:
    match source.type:
        case None:
            pass
        case MSO_COLOR_TYPE.RGB:
            copy_scalar_attributes(source, target, attribute_names=["rgb", "brightness"])
        case MSO_COLOR_TYPE.SCHEME:
            copy_scalar_attributes(source, target, attribute_names=["theme_color", "brightness"])
        case _:
            raise NotImplementedError(source.type)


def copy_data_label_properties(source: DataLabel, target: DataLabel) -> None:
    copy_scalar_attributes(source, target, attribute_names={"has_text_frame", "position"})
    if source.has_text_frame:
        copy_font_properties(cast(Font, source.font), cast(Font, target.font))
        copy_text_frame_properties(source.text_frame, target.text_frame)


def copy_data_labels_properties(source: DataLabels, target: DataLabels) -> None:
    copy_scalar_attributes(
        source,
        target,
        attribute_names={
            "number_format_is_linked",
            "number_format",
            "show_category_name",
            "show_legend_key",
            "show_percentage",
            "show_series_name",
            "show_value",
        },
    )
    copy_font_properties(cast(Font, source.font), cast(Font, target.font))


def copy_fill_format_properties(source: FillFormat, target: FillFormat) -> None:
    if source.type != target.type:
        match source.type:
            case MSO_FILL.BACKGROUND:
                target.background()
            case MSO_FILL.SOLID:
                target.solid()
            case _:
                raise NotImplementedError(source.type)
    match target.type:
        case None:
            pass
        case MSO_FILL.SOLID:
            copy_color_format_properties(cast(ColorFormat, source.fore_color), cast(ColorFormat, target.fore_color))
        case MSO_FILL.BACKGROUND:
            pass  # BACKGROUND means transparent
        case _:
            raise NotImplementedError(source.type)


def copy_font_properties(source: Font, target: Font) -> None:
    copy_scalar_attributes(
        source, target, attribute_names={"bold", "italic", "language_id", "name", "size", "underline"}
    )
    copy_color_format_properties(cast(ColorFormat, source.color), cast(ColorFormat, target.color))
    copy_fill_format_properties(cast(FillFormat, source.fill), cast(FillFormat, target.fill))


def copy_legend_properties(source: Legend, target: Legend) -> None:
    copy_scalar_attributes(source, target, attribute_names={"horz_offset", "include_in_layout", "position"})
    copy_font_properties(cast(Font, source.font), cast(Font, target.font))


def copy_line_format_properties(source: LineFormat, target: LineFormat) -> None:
    copy_scalar_attributes(source, target, attribute_names={"dash_style", "width"})
    copy_color_format_properties(cast(ColorFormat, source.color), cast(ColorFormat, target.color))
    copy_fill_format_properties(cast(FillFormat, source.fill), cast(FillFormat, target.fill))


def copy_major_gridlines_properties(source: MajorGridlines, target: MajorGridlines) -> None:
    copy_chart_format_properties(cast(ChartFormat, source.format), cast(ChartFormat, target.format))


def copy_plots_properties(source: _Plots, target: _Plots) -> None:
    for source_plot, target_plot in zip(
        cast(Iterable[_BasePlot], source), cast(Iterable[_BasePlot], target), strict=True
    ):
        copy_plot_properties(source_plot, target_plot)


def copy_plot_properties(source: _BasePlot, target: _BasePlot) -> None:
    try:
        copy_scalar_attributes(source, target, attribute_names={"has_data_labels", "vary_by_categories"})
        if source.has_data_labels:
            copy_data_labels_properties(source.data_labels, target.data_labels)
    except AttributeError:
        # XY charts series can't read has_data_labels: AttributeError: 'CT_ScatterChart' object has no attribute 'dLbls'
        pass


def copy_point_properties(source: Point, target: Point) -> None:
    copy_data_label_properties(cast(DataLabel, source.data_label), cast(DataLabel, target.data_label))
    copy_chart_format_properties(cast(ChartFormat, source.format), cast(ChartFormat, target.format))


def copy_scalar_attributes(source: object, target: object, *, attribute_names: Iterable[str]) -> None:
    for attribute_name in attribute_names:
        source_attribute_value = getattr(source, attribute_name)
        setattr(target, attribute_name, source_attribute_value)


def copy_series_collection_properties(source: SeriesCollection, target: SeriesCollection) -> None:
    for source_series, target_series in zip(source, target, strict=True):
        if isinstance(source_series, _BaseCategorySeries) and isinstance(target_series, _BaseCategorySeries):
            copy_base_category_series_properties(source_series, target_series)
        elif isinstance(source_series, XySeries) and isinstance(target_series, XySeries):
            copy_xy_series_properties(source_series, target_series)
        else:
            raise NotImplementedError


def copy_shape_properties(source: BaseShape, target: BaseShape) -> None:
    copy_scalar_attributes(source, target, attribute_names={"name", "rotation"})
    if source.has_chart and target.has_chart:
        assert isinstance(source, GraphicFrame)
        assert isinstance(target, GraphicFrame)
        copy_chart_properties(source.chart, target.chart)


def copy_text_frame_properties(source: TextFrame, target: TextFrame) -> None:
    copy_scalar_attributes(
        source,
        target,
        attribute_names={
            "margin_bottom",
            "margin_left",
            "margin_right",
            "margin_top",
            "text",
            "vertical_anchor",
            "word_wrap",
        },
    )


def copy_tick_labels_properties(source: TickLabels, target: TickLabels, *, axis: _BaseAxis) -> None:
    copy_scalar_attributes(source, target, attribute_names={"number_format", "number_format_is_linked"})
    copy_font_properties(cast(Font, source.font), cast(Font, target.font))
    if isinstance(axis, CategoryAxis):
        copy_scalar_attributes(source, target, attribute_names={"offset"})


def copy_value_axis_properties(source: ValueAxis, target: ValueAxis) -> None:
    """Copy only the properties that are specific to ValueAxis.

    Cf copy_axis_properties.
    """
    copy_scalar_attributes(
        source,
        target,
        attribute_names={"crosses", "crosses_at", "major_unit", "minor_unit"},
    )


def copy_xy_series_properties(source: XySeries, target: XySeries) -> None:
    # Ignore data_labels because they will be generated automatically
    copy_chart_format_properties(cast(ChartFormat, source.format), cast(ChartFormat, target.format))
