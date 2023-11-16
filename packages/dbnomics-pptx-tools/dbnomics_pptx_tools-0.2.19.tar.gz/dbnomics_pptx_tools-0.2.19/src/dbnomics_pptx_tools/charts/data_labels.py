import dataclasses
from collections.abc import Iterable, Iterator
from dataclasses import dataclass
from typing import Any, cast

import daiquiri
from lxml import etree
from pandas import DataFrame
from pptx.chart.chart import Chart
from pptx.chart.datalabel import DataLabel
from pptx.chart.plot import _BasePlot
from pptx.chart.point import Point
from pptx.chart.series import CategoryPoints, _BaseCategorySeries
from pptx.dml.line import LineFormat
from pptx.enum.chart import XL_DATA_LABEL_POSITION
from pptx.oxml.ns import nsdecls
from pptx.text.text import Font

from dbnomics_pptx_tools.charts.utils import iter_chart_series_with_spec
from dbnomics_pptx_tools.metadata import ChartSpec, DataLabelPosition, DataLabelShow, DataLabelSpec
from dbnomics_pptx_tools.pptx_copy import copy_color_format_properties, copy_font_properties
from dbnomics_pptx_tools.xml_utils import remove_element

logger = daiquiri.getLogger(__name__)


@dataclass
class DataLabelRenderData:
    chart_series: _BaseCategorySeries
    point: Point
    ratio: float
    new_ratio: float | None

    @property
    def ratio_distance(self) -> float | None:
        if self.new_ratio is None:
            return None
        return self.new_ratio - self.ratio


def add_data_labels_to_each_series(
    chart: Chart,
    *,
    chart_spec: ChartSpec,
    data_label_spec: DataLabelSpec,
    pivoted_df: DataFrame,
    points_by_series_name: dict[str, list[Point]],
) -> None:
    if data_label_spec.position == DataLabelPosition.NONE:
        return

    if data_label_spec.position == DataLabelPosition.ALL_POINTS:
        add_data_labels_to_every_point_of_each_series(
            chart, chart_spec=chart_spec, data_label_spec=data_label_spec, pivoted_df=pivoted_df
        )
    elif data_label_spec.position == DataLabelPosition.FIRST_AND_LAST_POINTS:
        add_data_labels_to_first_and_last_points_of_each_series(
            chart, data_label_spec=data_label_spec, pivoted_df=pivoted_df
        )
    elif data_label_spec.position == DataLabelPosition.LAST_POINT:
        add_data_labels_to_last_point_of_each_series(
            chart,
            chart_spec=chart_spec,
            data_label_spec=data_label_spec,
            pivoted_df=pivoted_df,
            points_by_series_name=points_by_series_name,
        )


def add_data_labels_to_every_point_of_each_series(
    chart: Chart, *, chart_spec: ChartSpec, data_label_spec: DataLabelSpec, pivoted_df: DataFrame
) -> None:
    logger.debug("Adding a data label to every point of each series of the chart...")

    value_axis_font = cast(Font, chart.value_axis.tick_labels.font)
    positions = [XL_DATA_LABEL_POSITION.ABOVE, XL_DATA_LABEL_POSITION.BELOW]

    for index, (chart_series, series_spec) in enumerate(
        iter_chart_series_with_spec(chart, chart_series_attrs=pivoted_df.attrs, chart_spec=chart_spec)
    ):
        position = positions[index % len(positions)]
        for point, (period, value) in zip(
            cast(CategoryPoints, chart_series.points), pivoted_df[series_spec.id].items(), strict=True
        ):
            add_data_label_to_point(
                point,
                text=get_shown(data_label_spec, period, value),
                font=value_axis_font,
                number_format=data_label_spec.number_format,
                position=position,
            )


def add_data_label_to_point(
    point: Point,
    *,
    text: str,
    font: Font | None = None,
    number_format: str,
    position: XL_DATA_LABEL_POSITION = XL_DATA_LABEL_POSITION.ABOVE,
) -> None:
    data_label = cast(DataLabel, point.data_label)
    data_label.has_text_frame = True
    data_label.text_frame.text = text
    data_label.position = position
    if font is not None:
        for paragraph in data_label.text_frame.paragraphs:
            copy_font_properties(font, paragraph.font)
    add_number_format(data_label, number_format)


def add_data_labels_to_first_and_last_points_of_each_series(
    chart: Chart, *, data_label_spec: DataLabelSpec, pivoted_df: DataFrame
) -> None:
    logger.debug("Adding a data label to the first and last points of each series of the chart...")

    show = data_label_spec.show
    value_axis_font = cast(Font, chart.value_axis.tick_labels.font)

    for chart_series in chart.series:
        if show is not DataLabelShow.PERIOD:
            raise NotImplementedError(show)

        points = cast(CategoryPoints, chart_series.points)

        add_data_label_to_point(
            points[0],
            text=pivoted_df.index[0],
            font=value_axis_font,
            number_format=data_label_spec.number_format,
        )
        add_data_label_to_point(
            points[len(points) - 1],
            text=pivoted_df.index[-1],
            font=value_axis_font,
            number_format=data_label_spec.number_format,
        )


def add_data_labels_to_last_point_of_each_series(
    chart: Chart,
    *,
    chart_spec: ChartSpec,
    data_label_spec: DataLabelSpec,
    pivoted_df: DataFrame,
    points_by_series_name: dict[str, list[Point]],
) -> None:
    logger.debug("Adding a data label to the last point of each series of the chart...")

    render_data_list = compute_data_label_positions(
        chart, chart_spec=chart_spec, pivoted_df=pivoted_df, points_by_series_name=points_by_series_name
    )

    for render_data in render_data_list:
        apply_render_data_to_chart(render_data, chart=chart, data_label_spec=data_label_spec)


def add_number_format(data_label: DataLabel, number_format: str) -> None:
    d_lbl_element = data_label._get_or_add_dLbl()  # noqa: SLF001
    num_fmt_element = d_lbl_element.find("./{*}numFmt")
    if num_fmt_element is None:
        num_fmt_element = etree.fromstring(f"""<c:numFmt {nsdecls("c")} />""")  # noqa: S320
        d_lbl_element.append(num_fmt_element)
    num_fmt_element.attrib["formatCode"] = number_format
    num_fmt_element.attrib["sourceLinked"] = "0"
    d_lbl_element.find("./{*}showVal").attrib["val"] = "1"


def apply_render_data_to_chart(
    render_data: DataLabelRenderData, *, chart: Chart, data_label_spec: DataLabelSpec
) -> None:
    logger.debug(
        "Adding data label to chart for series %r (using number format %r)...",
        render_data.chart_series.name,
        data_label_spec.number_format,
    )
    value_axis_font = cast(Font, chart.value_axis.tick_labels.font)

    data_label = cast(DataLabel, render_data.point.data_label)
    copy_font_properties(value_axis_font, data_label.font)

    add_number_format(data_label, data_label_spec.number_format)

    line = LineFormat(data_label._dLbl.get_or_add_spPr())  # noqa: SLF001
    copy_color_format_properties(render_data.chart_series.format.line.color, line.color)

    ratio_distance = render_data.ratio_distance
    if ratio_distance is not None:
        logger.debug(
            "Moving the data label of the series %r because if is too close to the previous one",
            render_data.chart_series.name,
        )
        layout_element = etree.fromstring(  # noqa: S320
            f"""
                <c:layout {nsdecls("c")}>
                    <c:manualLayout>
                        <c:x val="0"/>
                        <c:y val="{-ratio_distance}"/>
                    </c:manualLayout>
                </c:layout>
            """.strip()
        )
        d_lbl_element = data_label._get_or_add_dLbl()  # noqa: SLF001
        d_lbl_element.append(layout_element)


def compute_data_label_positions(
    chart: Chart, *, chart_spec: ChartSpec, pivoted_df: DataFrame, points_by_series_name: dict[str, list[Point]]
) -> list[DataLabelRenderData]:
    logger.debug("Computing data label positions...")
    render_data_list: list[DataLabelRenderData] = []

    for chart_series, series_spec in iter_chart_series_with_spec(
        chart, chart_series_attrs=pivoted_df.attrs, chart_spec=chart_spec
    ):
        series_id = series_spec.id
        series_attrs = pivoted_df.attrs.get(series_id, {})
        series_name = series_spec.name.format(**series_attrs)

        series = pivoted_df.reset_index()[series_id]
        last_value_index = cast(int, series.last_valid_index())
        if last_value_index is None:
            logger.warning("The series named %r with ID %r has only NA values, skipping", series_name, series_id)
            continue

        last_value = series[last_value_index]
        ratio = compute_data_label_ratio(last_value, chart=chart, pivoted_df=pivoted_df)
        last_point = points_by_series_name[series_name][last_value_index]
        render_data_list.append(
            DataLabelRenderData(chart_series=chart_series, point=last_point, ratio=ratio, new_ratio=None)
        )

    render_data_list = sorted(render_data_list, key=lambda render_data: render_data.ratio)
    return list(iter_spaced_data_labels(render_data_list))


def compute_data_label_ratio(value: float, *, chart: Chart, pivoted_df: DataFrame) -> float:
    chart_min_value, chart_max_value = compute_value_axis_bounds(pivoted_df, chart=chart)
    chart_value_range = chart_max_value - chart_min_value
    return (value - chart_min_value) / chart_value_range


def compute_value_axis_bounds(pivoted_df: DataFrame, *, chart: Chart, margin_ratio: float = 0.1) -> tuple[float, float]:
    min_value = pivoted_df.min().min()
    max_value = pivoted_df.max().max()
    margin = (max_value - min_value) * margin_ratio
    minimum_scale = chart.value_axis.minimum_scale
    maximum_scale = chart.value_axis.maximum_scale
    return (
        minimum_scale if minimum_scale is not None else min_value - margin,
        maximum_scale if maximum_scale is not None else max_value + margin,
    )


def get_shown(data_label_spec: DataLabelSpec, period: Any, value: Any) -> str:
    show = data_label_spec.show
    if show is DataLabelShow.PERIOD:
        return str(period)
    if show is DataLabelShow.VALUE:
        num_decimals = len(data_label_spec.number_format.split(".")[1])
        format_str = f"{{:.{num_decimals}f}}"
        return format_str.format(value)
    raise NotImplementedError(show)


def iter_spaced_data_labels(
    render_data_list: list[DataLabelRenderData], *, min_ratio_distance: float = 0.05
) -> Iterator[DataLabelRenderData]:
    if not render_data_list:
        return []

    yield render_data_list[0]
    last_ratio = render_data_list[0].ratio

    for current in render_data_list[1:]:
        if current.ratio - last_ratio < min_ratio_distance:
            new_ratio = last_ratio + min_ratio_distance
            yield dataclasses.replace(current, new_ratio=new_ratio)
            last_ratio = new_ratio
        else:
            yield current
            last_ratio = current.ratio


def remove_data_labels(chart: Chart, *, points_by_series_name: dict[str, list[Point]]) -> None:
    for plot_index, plot in enumerate(cast(Iterable[_BasePlot], chart.plots)):
        try:
            if plot.has_data_labels:
                logger.debug("Plot #%d has data labels, removing", plot_index)
                plot.has_data_labels = False
        except Exception:
            logger.exception("Could not do plot.has_data_labels = False")
        for series in cast(Iterable[_BaseCategorySeries], plot.series):
            series_points = points_by_series_name[series.name]
            for point_position, point in enumerate(series_points, start=1):
                data_label = cast(DataLabel, point.data_label)
                d_lbl = data_label._dLbl  # noqa: SLF001
                if d_lbl is not None:
                    logger.debug(
                        "Point %d/%d of series %r has a data label, removing",
                        point_position,
                        len(series_points),
                        series.name,
                    )
                    remove_element(d_lbl)


def update_data_labels(chart: Chart, *, chart_spec: ChartSpec, pivoted_df: DataFrame) -> None:
    points_by_series_name = {series.name: list(series.points) for series in chart.series}
    remove_data_labels(chart, points_by_series_name=points_by_series_name)
    add_data_labels_to_each_series(
        chart,
        chart_spec=chart_spec,
        data_label_spec=chart_spec.get_data_label_spec(),
        pivoted_df=pivoted_df,
        points_by_series_name=points_by_series_name,
    )
