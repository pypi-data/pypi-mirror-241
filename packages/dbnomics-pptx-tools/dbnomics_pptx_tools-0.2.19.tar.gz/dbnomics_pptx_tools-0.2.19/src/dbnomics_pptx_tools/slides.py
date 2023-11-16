from collections.abc import Container, Iterable, Iterator
from typing import TYPE_CHECKING, cast

import daiquiri
from more_itertools import first_true
from parsy import ParseError, any_char, regex, string, whitespace
from pptx.chart.series import LineSeries, SeriesCollection
from pptx.presentation import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide, Slides

from dbnomics_pptx_tools.charts import update_chart
from dbnomics_pptx_tools.repo import SeriesRepo
from dbnomics_pptx_tools.tables import extract_table_zones, format_table, update_table

from .metadata import PresentationMetadata, SlideMetadata

if TYPE_CHECKING:
    from pptx.table import Table

logger = daiquiri.getLogger(__name__)


def delete_other_slides(prs: Presentation, *, only_slides: Container[int]) -> None:
    for slide_pos, slide in enumerate(cast(Iterable[Slide], prs.slides), start=1):
        if slide_pos not in only_slides:
            delete_slide(prs, slide)


# Cf https://github.com/scanny/python-pptx/issues/67#issuecomment-296135015
def delete_slide(prs: Presentation, slide: Slide) -> None:
    slides = cast(Slides, prs.slides)
    sld_id_lst = slides._sldIdLst  # noqa: SLF001
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(sld_id_lst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del sld_id_lst[id_dict[slide_id][0]]


def extract_slide_id_from_slide_notes(slide: Slide) -> str | None:
    slide_notes = slide.notes_slide.notes_text_frame.text
    parser = (
        whitespace.optional()
        >> string("slide_id", transform=str.lower)
        >> whitespace.optional()
        >> string(":")
        >> whitespace.optional()
        >> regex(r"[\w_-]+")
        << any_char.many().optional()
    )
    try:
        return cast(str, parser.parse(slide_notes))
    except ParseError:
        return None


def find_slide_by_id(prs: Presentation, *, slide_id: str) -> Slide | None:
    for slide in cast(Iterable[Slide], prs.slides):
        current_slide_id = extract_slide_id_from_slide_notes(slide)
        if current_slide_id is not None and slide_id == current_slide_id:
            return slide

    return None


def find_table_shape(slide: Slide, table_name: str) -> GraphicFrame | None:
    return first_true(
        iter_table_shapes(slide),
        default=None,
        pred=lambda table_shape: table_shape.name == table_name,  # type: ignore[union-attr]
    )


def iter_chart_shapes(slide: Slide) -> Iterator[GraphicFrame]:
    for shape in cast(list[GraphicFrame], slide.shapes):
        if shape.has_chart:
            yield shape


def iter_table_shapes(slide: Slide) -> Iterator[GraphicFrame]:
    for shape in cast(list[GraphicFrame], slide.shapes):
        if shape.has_table:
            yield shape


def iter_line_series(series_collection: SeriesCollection, *, strict: bool = True) -> Iterator[LineSeries]:
    for series in series_collection:
        if not isinstance(series, LineSeries):
            if strict:
                msg = f"Only LineSeries are expected in {series_collection!r}"
                raise ValueError(msg)
            else:
                continue
        yield series


def update_slide_charts(
    slide: Slide, *, fail_fast: bool = False, repo: SeriesRepo, slide_metadata: SlideMetadata
) -> None:
    slide_shapes = list(iter_chart_shapes(slide))
    for chart_pos, chart_shape in enumerate(slide_shapes, start=1):
        chart_name = chart_shape.name
        chart_full_name = f"{chart_name} ({chart_pos}/{len(slide_shapes)})"

        chart_spec = slide_metadata.charts.get(chart_name)
        if chart_spec is None:
            logger.debug(
                "Chart %s was found in the presentation file, "
                "but no corresponding block was found in the YAML file, ignoring chart",
                chart_full_name,
            )
            continue

        logger.debug("Updating chart %s...", chart_full_name)

        try:
            update_chart(
                chart_shape,
                chart_spec=chart_spec,
                repo=repo,
                slide=slide,
            )
        except Exception:
            if fail_fast:
                raise
            logger.exception("Error updating chart %s, skipping", chart_full_name)


def update_slide_tables(
    slide: Slide, *, fail_fast: bool = False, repo: SeriesRepo, slide_metadata: SlideMetadata
) -> None:
    table_shapes = list(iter_table_shapes(slide))
    for table_pos, table_shape in enumerate(table_shapes, start=1):
        table: Table = table_shape.table
        table_name = table_shape.name
        table_full_name = f"{table_name} ({table_pos}/{len(table_shapes)})"

        table_spec = slide_metadata.tables.get(table_name)
        if not table_spec:
            logger.debug(
                "Table %s was found in the presentation file, "
                "but no corresponding block was found in the YAML file, ignoring table:\n%s",
                table_full_name,
                format_table(table),
            )
            continue

        logger.debug("Updating table %s...", table_full_name)

        table_zones = extract_table_zones(table, table_spec=table_spec)
        if table_zones is None:
            logger.warning(
                "Could not extract the zones of table %s, ignoring table:\n%s",
                table_full_name,
                format_table(table),
            )
            continue

        try:
            update_table(table, repo=repo, table_spec=table_spec, table_zones=table_zones)
        except Exception:
            if fail_fast:
                raise
            logger.exception("Error updating table %s, skipping", table_full_name)


def update_slides(
    prs: Presentation,
    *,
    fail_fast: bool = False,
    only_slides: Container[int] | None = None,
    presentation_metadata: PresentationMetadata,
    repo: SeriesRepo,
) -> None:
    for slide_pos, slide in enumerate(cast(Iterable[Slide], prs.slides), start=1):
        if only_slides is not None and slide_pos not in only_slides:
            continue

        slide_title: str = slide.shapes.title.text

        slide_id = extract_slide_id_from_slide_notes(slide)
        if slide_id is None:
            logger.debug("Could not extract the slide ID from the slide notes, using the slide title %r", slide_title)

        slide_metadata = presentation_metadata.slides.get(slide_id if slide_id is not None else slide_title)
        if slide_metadata is None:
            logger.debug(
                "Slide %d with ID %r and title %r has no metadata, ignoring slide", slide_pos, slide_id, slide_title
            )
            continue

        logger.debug("Updating slide %d with ID %r and title %r...", slide_pos, slide_id, slide_title)

        update_slide_charts(slide, fail_fast=fail_fast, repo=repo, slide_metadata=slide_metadata)

        update_slide_tables(slide, fail_fast=fail_fast, repo=repo, slide_metadata=slide_metadata)
