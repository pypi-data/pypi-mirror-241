import logging
from collections.abc import Iterable
from pathlib import Path
from typing import Final, Optional, cast

import daiquiri
import typer
from pptx.slide import Slide
from typer import FileBinaryRead, FileBinaryWrite, Typer

from dbnomics_pptx_tools import tables
from dbnomics_pptx_tools.cache import SeriesCache
from dbnomics_pptx_tools.cli.cli_utils import (
    load_presentation_metadata,
    open_presentation,
    parse_slide_option,
    parse_slides_option,
)
from dbnomics_pptx_tools.data_archive import save_data_archive_file
from dbnomics_pptx_tools.metadata import PresentationMetadata
from dbnomics_pptx_tools.repo import SeriesRepo, SeriesRepoError
from dbnomics_pptx_tools.slides import (
    delete_other_slides,
    extract_slide_id_from_slide_notes,
    find_slide_by_id,
    find_table_shape,
    update_slides,
)
from dbnomics_pptx_tools.tables import format_table

app = Typer()

logger = daiquiri.getLogger(__name__)


DBNOMICS_API_CACHE_DIR_NAME: Final = "dbnomics_api_cache"
DEFAULT_TIMEOUT: Final = 60


@app.callback(context_settings={"help_option_names": ["-h", "--help"]})
def main(verbose: bool = typer.Option(False, "-v", help="Show debug log messages")) -> None:
    """DBnomics PowerPoint (pptx) tools."""
    daiquiri.setup()
    if verbose:
        daiquiri.set_default_log_levels([("dbnomics_pptx_tools", logging.DEBUG)])


@app.command()
def extract_table_zones(
    input_pptx_file: FileBinaryRead,
    slide_expr: str = typer.Argument(..., help="Slide ID or number"),
    table_name: str = typer.Argument(...),
    metadata_file: Path = typer.Option(..., exists=True, readable=True),
) -> None:
    prs = open_presentation(input_pptx_file)
    prs_slide_ids = [extract_slide_id_from_slide_notes(slide) for slide in cast(Iterable[Slide], prs.slides)]
    slide_number = parse_slide_option(slide_expr, slide_ids=prs_slide_ids)
    slide_id = prs_slide_ids[slide_number - 1]
    if slide_id is None:
        msg = f"Slide number {slide_number} does not have an ID defined in the slide notes"
        raise typer.BadParameter(msg)

    presentation_metadata = load_presentation_metadata(metadata_file)
    slide_metadata = presentation_metadata.slides.get(slide_id)
    table_spec = None if slide_metadata is None else slide_metadata.tables.get(table_name)

    slide = find_slide_by_id(prs, slide_id=slide_id)
    if slide is None:
        typer.echo(f"Could not find slide wiht ID {slide_id!r}")
        raise typer.Exit(1)
    logger.debug("Found slide %r", slide)

    table_shape = find_table_shape(slide, table_name)
    if table_shape is None:
        typer.echo(f"Could not find table {table_name!r} in slide with ID {slide_id!r}")
        raise typer.Exit(1)
    logger.debug("Found table shape %r", table_shape)

    table = table_shape.table
    logger.debug("Showing table preview:\n%s", format_table(table))

    table_zones = tables.extract_table_zones(table, table_spec=table_spec)
    if table_zones is None:
        typer.echo(f"Could not extract the zones of table {table_name!r}")
        raise typer.Exit(1)

    typer.echo(table_zones)


@app.command()
def fetch(
    metadata_file: Path = typer.Argument(..., exists=True, readable=True),
    dbnomics_api_cache_dir: Path = typer.Option(DBNOMICS_API_CACHE_DIR_NAME),
    resume: bool = typer.Option(False, help="Do not fetch the series that are already stored in the cache."),
    timeout: int = typer.Option(DEFAULT_TIMEOUT, help="Timeout in seconds of every HTTP request to DBnomics API"),
) -> None:
    presentation_metadata = load_presentation_metadata(metadata_file)
    fetchable_series_ids = sorted(presentation_metadata.find_fetchable_series_ids())

    cache = SeriesCache(cache_dir=dbnomics_api_cache_dir)
    repo = SeriesRepo(auto_fetch=True, cache=cache, resume=resume, timeout=timeout)

    cached_series_ids = []
    series_ids_to_fetch = fetchable_series_ids
    if resume:
        cached_series_ids = [series_id for series_id in fetchable_series_ids if cache.has_df(series_id)]
        logger.debug(
            "The following %d series won't be fetched because they are already stored in the cache: %r",
            len(cached_series_ids),
            cached_series_ids,
        )
        series_ids_to_fetch = sorted(set(fetchable_series_ids) - set(cached_series_ids))

    failed_series_ids = []
    successful_series_ids = []
    if series_ids_to_fetch:
        logger.debug(
            "About to fetch the following %d series needed for the presentation: %r",
            len(series_ids_to_fetch),
            series_ids_to_fetch,
        )

        for series_id in series_ids_to_fetch:
            try:
                repo.fetch_csv(series_id)
                repo.load_df(series_id)
            except Exception:
                logger.exception("Error loading series %r, skipping", series_id)
                failed_series_ids.append(series_id)
            else:
                successful_series_ids.append(series_id)

    logger.info(
        "Attempted to fetch %d series, among them %d were successful, %d failed, %d were skipped because already in cache",  # noqa: E501
        len(fetchable_series_ids),
        len(successful_series_ids),
        len(failed_series_ids),
        len(cached_series_ids),
    )


@app.command()
def save_data_archive(
    metadata_file: Path = typer.Argument(..., exists=True, readable=True),
    output_data_archive_file: Path = typer.Argument(
        ..., help="Write a ZIP file containing data used to update the presentation", writable=True
    ),
    auto_fetch: bool = typer.Option(
        True, help="Fetch series when it is not found in the cache, then add it to the cache."
    ),
    dbnomics_api_cache_dir: Path = typer.Option(DBNOMICS_API_CACHE_DIR_NAME),
    force: bool = typer.Option(False, help="Fetch a series even if it is already stored in the cache."),
    timeout: int = typer.Option(DEFAULT_TIMEOUT, help="Timeout in seconds of every HTTP request to DBnomics API"),
) -> None:
    """Save a data ZIP archive of the series used to update the presentation."""
    if not output_data_archive_file.name.endswith(".zip"):
        msg = f"Only ZIP archives are supported: the file name must end with '.zip', got {output_data_archive_file.name!r}"  # noqa: E501
        raise typer.BadParameter(msg)

    presentation_metadata = load_presentation_metadata(metadata_file)

    cache = SeriesCache(cache_dir=dbnomics_api_cache_dir)
    repo = SeriesRepo(auto_fetch=auto_fetch, cache=cache, resume=not force, timeout=timeout)

    if output_data_archive_file is not None:
        save_data_archive_file(
            output_file=output_data_archive_file, presentation_metadata=presentation_metadata, repo=repo
        )


@app.command()
def show_presentation_metadata_json_schema() -> None:
    """Show JSON schema of presentation metadata."""
    typer.echo(PresentationMetadata.schema_json(indent=2))


@app.command()
def update(
    input_pptx_file: FileBinaryRead,
    output_pptx_file: FileBinaryWrite,
    auto_fetch: bool = typer.Option(
        True, help="Fetch series when it is not found in the cache, then add it to the cache."
    ),
    dbnomics_api_cache_dir: Path = typer.Option(DBNOMICS_API_CACHE_DIR_NAME),
    fail_fast: bool = typer.Option(False, help="Raise exception when failing to update a slide."),
    force: bool = typer.Option(False, help="Fetch a series even if it is already stored in the cache."),
    metadata_file: Path = typer.Option(..., exists=True, readable=True),
    only_slides_expr: Optional[str] = typer.Option(None, "--slides"),
    save_processed_slides_only: bool = typer.Option(False),
    timeout: int = typer.Option(DEFAULT_TIMEOUT, help="Timeout in seconds of every HTTP request to DBnomics API"),
) -> None:
    """Update DBnomics data in a PowerPoint (pptx) presentation."""
    prs = open_presentation(input_pptx_file)

    only_slides = None
    if only_slides_expr is not None:
        logger.debug("Will process slides %s", only_slides_expr)
        prs_slide_ids = [extract_slide_id_from_slide_notes(slide) for slide in cast(Iterable[Slide], prs.slides)]
        only_slides = parse_slides_option(only_slides_expr, slide_ids=prs_slide_ids)

    if save_processed_slides_only and only_slides is None:
        msg = "--save-processed-slides-only must be used with --slides"
        raise typer.BadParameter(msg)

    presentation_metadata = load_presentation_metadata(metadata_file)

    cache = SeriesCache(cache_dir=dbnomics_api_cache_dir)
    repo = SeriesRepo(auto_fetch=auto_fetch, cache=cache, resume=not force, timeout=timeout)

    try:
        update_slides(
            prs, fail_fast=fail_fast, only_slides=only_slides, presentation_metadata=presentation_metadata, repo=repo
        )
    except SeriesRepoError as exc:
        typer.echo(f'{exc!s} Hint: use the --auto-fetch option or run the "fetch" command first.')
        raise typer.Exit(1) from exc

    if save_processed_slides_only:
        assert only_slides is not None
        delete_other_slides(prs, only_slides=only_slides)

    prs.save(output_pptx_file)
    logger.info("Output presentation was saved as %r", str(output_pptx_file.name))


if __name__ == "__main__":
    app()
